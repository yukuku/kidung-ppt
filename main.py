import codecs
import sys

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE

import songbook_parser


def main():
    if len(sys.argv) < 3:
        print("usage: {} <songbook file> <ppt file>".format(sys.argv[0]))
        return

    sbf = sys.argv[1]
    ppf = sys.argv[2]

    print('Processing {}...'.format(sbf))

    prs = Presentation('song_template.pptx')

    def new_slide(title_text, paras, code_text):
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # for shape in slide.placeholders:
        #     print('%d (%s): %s' % (shape.placeholder_format.idx, shape.placeholder_format.type, shape.name))

        title = slide.placeholders[0]
        body = slide.placeholders[1]
        codepl = slide.placeholders[13]

        title.text = title_text
        codepl.text = code_text

        tf = body.text_frame

        for i, para in enumerate(paras):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = para

        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    songs = songbook_parser.parse_book_content(codecs.open(sbf, encoding='utf8').read())
    for song in songs:
        print('Song {} {}...'.format(song['code'], song['title']))
        for lyric in song['lyrics']:
            for verse in lyric['verses']:
                nonrefversecount = sum(v['kind'] != 'VerseKind.REFRAIN' for v in lyric['verses'])
                if verse['kind'] == 'VerseKind.REFRAIN':
                    title_text = '{} (ref)'.format(song['title'], verse['ordering'])
                elif nonrefversecount == 1:
                    title_text = song['title']
                else:
                    title_text = '{} ({})'.format(song['title'], verse['ordering'])

                new_slide(title_text, verse['lines'], 'KPRI {}'.format(song['code']))

    prs.save(ppf)


if __name__ == '__main__':
    main()
